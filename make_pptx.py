"""
make_pptx.py — Génère la présentation PowerPoint du Projet C
Usage : python make_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ───────────────────────────────────────────────────
GREEN      = RGBColor(0x27, 0xAE, 0x60)   # accent principal
DARK       = RGBColor(0x1A, 0x1A, 0x2E)   # fond sombre
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4)   # fond clair
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT  = RGBColor(0x55, 0x55, 0x55)
CODE_BG    = RGBColor(0x2D, 0x2D, 0x2D)
CODE_FG    = RGBColor(0xA8, 0xFF, 0x78)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout vide


# ── Helpers ───────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, color, radius=False):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, items, l, t, w, h, size=17, color=DARK, bullet="▸ "):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(6)

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, DARK)
    add_text(slide, title, 0.4, 0.15, 12, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12, 0.4, size=14, color=GREEN)
    add_rect(slide, 0, 1.3, 13.33, 0.04, GREEN)

def code_block(slide, code, l, t, w, h, size=11):
    add_rect(slide, l, t, w, h, CODE_BG)
    txBox = slide.shapes.add_textbox(Inches(l+0.15), Inches(t+0.15),
                                      Inches(w-0.3), Inches(h-0.3))
    tf = txBox.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = CODE_FG
        run.font.name = "Courier New"


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Titre
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, DARK)
add_rect(s, 0, 3.3, 13.33, 0.06, GREEN)

add_text(s, "Système de suivi des stocks", 1, 1.2, 11.33, 1.2,
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "TP MongoDB — Optimisation BDD NoSQL", 1, 2.5, 11.33, 0.6,
         size=20, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s, "Projet C", 1, 3.1, 11.33, 0.5,
         size=16, color=GRAY_TEXT, align=PP_ALIGN.CENTER)
add_text(s, "Kelvya  •  Maxime", 1, 4.2, 11.33, 0.6,
         size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "2025 – 2026", 1, 4.9, 11.33, 0.5,
         size=14, color=GRAY_TEXT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Présentation du sujet
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "Présentation du sujet", "Contexte & données")

add_text(s, "Contexte", 0.5, 1.6, 5.5, 0.4, size=16, bold=True, color=DARK)
add_bullet_box(s, [
    "Gérer les stocks d'une entreprise",
    "Suivre les fournisseurs et leurs délais",
    "Surveiller les commandes et leur statut",
], 0.5, 2.0, 5.8, 2.0, color=DARK)

add_text(s, "4 fonctionnalités obligatoires", 7.0, 1.6, 5.8, 0.4, size=16, bold=True, color=DARK)
add_bullet_box(s, [
    "F1 — État du stock d'un produit",
    "F2 — Commandes bloquées (stock insuffisant)",
    "F3 — Délais de livraison par fournisseur",
    "F4 — Date de livraison d'une commande",
], 7.0, 2.0, 5.8, 2.5, color=DARK)

# Cards données
cards = [("50", "Produits"), ("20", "Fournisseurs"), ("450", "Commandes"), ("100", "Relations\nfournisseurs")]
for i, (nb, label) in enumerate(cards):
    x = 0.5 + i * 3.1
    add_rect(s, x, 4.8, 2.8, 1.8, DARK)
    add_text(s, nb, x, 5.0, 2.8, 0.7, size=32, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, label, x, 5.7, 2.8, 0.6, size=14, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Choix technologiques
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "Choix technologiques", "Stack utilisé")

techs = [
    ("MongoDB", "Base NoSQL orientée documents\nCP selon le théorème de CAP\nIdéale pour données hiérarchiques"),
    ("PyMongo", "Pilote natif Python\nAccès direct aux collections\nAgrégations & index"),
    ("Flask", "Framework web Python léger\nRoutes = fonctionnalités métier\nTemplates Jinja2"),
    ("Python 3", "Langage unique front + back\nScript d'import + app web\nFacile à démontrer"),
]
for i, (name, desc) in enumerate(techs):
    x = 0.4 + i * 3.2
    add_rect(s, x, 1.6, 2.9, 4.5, DARK)
    add_text(s, name, x, 1.7, 2.9, 0.6, size=18, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_rect(s, x+0.1, 2.35, 2.7, 0.04, GREEN)
    add_text(s, desc, x+0.1, 2.5, 2.7, 3.2, size=14, color=WHITE)

add_text(s, "Pourquoi MongoDB et pas SQL ?", 0.5, 6.1, 12, 0.4, size=14, bold=True, color=DARK)
add_text(s, "Documents JSON natifs · Pas de jointures coûteuses · Dénormalisation · Schéma flexible",
         0.5, 6.45, 12, 0.4, size=13, color=GRAY_TEXT)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture des collections
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "Architecture des collections", "4 collections MongoDB")

cols = [
    ("products", ["product_id", "name", "category", "unit_price", "unit", "► stock {}"]),
    ("suppliers", ["supplier_id", "name", "email", "phone", "country", "avg_delivery_days"]),
    ("supplier_products", ["sp_id", "supplier_id", "► supplier_name", "► supplier_country", "product_id", "delivery_days", "unit_price"]),
    ("orders", ["order_id", "product_id", "► product_name", "► product_unit", "quantity_ordered", "status", "order_date", "expected_delivery_date", "actual_delivery_date"]),
]

for i, (name, fields) in enumerate(cols):
    x = 0.3 + i * 3.25
    add_rect(s, x, 1.5, 3.0, 0.55, DARK)
    add_text(s, name, x, 1.52, 3.0, 0.5, size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_rect(s, x, 2.05, 3.0, 4.6, WHITE)
    for j, field in enumerate(fields):
        color = GREEN if field.startswith("►") else DARK
        add_text(s, field.replace("► ", ""), x+0.15, 2.1 + j*0.5, 2.7, 0.45, size=12,
                 color=color, bold=field.startswith("►"))

add_text(s, "► Champ dénormalisé (embarqué depuis une autre collection)", 0.5, 6.8, 12, 0.4,
         size=12, color=GREEN, italic=True)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Dénormalisation
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "Dénormalisation", "Choix architectural clé")

# Avant / Après
add_rect(s, 0.4, 1.5, 4.0, 0.5, RGBColor(0xE7, 0x4C, 0x3C))
add_text(s, "AVANT — Modèle relationnel", 0.4, 1.5, 4.0, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_bullet_box(s, ["products", "stock (table séparée)", "suppliers", "supplier_products", "orders"],
               0.5, 2.1, 3.8, 2.5, size=14, color=DARK, bullet="  • ")

add_rect(s, 5.0, 1.5, 4.0, 0.5, GREEN)
add_text(s, "APRÈS — Dénormalisé MongoDB", 5.0, 1.5, 4.0, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_bullet_box(s, ["products (+ stock embarqué)", "suppliers", "supplier_products (+ nom/pays)", "orders (+ snapshot produit)"],
               5.1, 2.1, 3.8, 2.5, size=14, color=DARK, bullet="  • ")

add_text(s, "→", 4.3, 2.8, 0.6, 0.5, size=28, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# Tableau
headers_row = ["Collection", "Données embarquées", "Bénéfice"]
rows = [
    ["products",          "stock {}",                   "0 jointure pour l'état du stock"],
    ["orders",            "product_name, product_unit", "Snapshot historique du produit"],
    ["supplier_products", "supplier_name, country",     "Évite $lookup sur suppliers"],
]
col_widths = [2.5, 3.5, 4.5]
col_x = [0.4, 2.95, 6.5]
row_y = [4.7, 5.2, 5.7, 6.2]

for ci, (hdr, cx, cw) in enumerate(zip(headers_row, col_x, col_widths)):
    add_rect(s, cx, row_y[0], cw, 0.45, DARK)
    add_text(s, hdr, cx+0.05, row_y[0]+0.05, cw-0.1, 0.35, size=13, bold=True, color=WHITE)

for ri, row in enumerate(rows):
    bg = WHITE if ri % 2 == 0 else RGBColor(0xEB, 0xF5, 0xEB)
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_widths)):
        add_rect(s, cx, row_y[ri+1], cw, 0.45, bg)
        color = GREEN if ci == 1 else DARK
        add_text(s, cell, cx+0.05, row_y[ri+1]+0.05, cw-0.1, 0.35, size=12, color=color)


# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Validation $jsonSchema
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "Validation des documents", "$jsonSchema sur la collection products")

add_bullet_box(s, [
    "Garantit la structure de chaque document inséré",
    "Erreur levée si un champ obligatoire manque ou a le mauvais type",
    "Appliqué à la collection products au moment de sa création",
], 0.5, 1.5, 6.0, 2.0, size=15, color=DARK)

code = '''db.create_collection("products", validator={
  "$jsonSchema": {
    "bsonType": "object",
    "required": ["product_id","name","category",
                 "unit_price","unit","stock"],
    "properties": {
      "product_id": { "bsonType": "int" },
      "category":   { "bsonType": "string",
                      "enum": ["Electronique","Alimentaire",
                               "Vetement","Outillage","Mobilier"] },
      "unit_price": { "bsonType": "double", "minimum": 0 },
      "stock": {
        "bsonType": "object",
        "required": ["quantity","min_threshold"],
        "properties": {
          "quantity": { "bsonType": "int", "minimum": 0 }
        }
      }
    }
  }
})'''
code_block(s, code, 6.8, 1.4, 6.1, 5.6, size=10)

add_text(s, "Règles définies", 0.5, 3.3, 5.5, 0.4, size=14, bold=True, color=DARK)
add_bullet_box(s, [
    "bsonType → type strict (int, string, double…)",
    "enum → valeurs autorisées pour category",
    "minimum → unit_price ≥ 0",
    "required → champs obligatoires dont stock {}",
], 0.5, 3.7, 5.8, 2.5, size=14, color=DARK)


# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Vues MongoDB
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "Vues MongoDB", "Collections en lecture seule basées sur des pipelines")

views = [
    ("vue_stock_status",          "Statut de chaque produit :\nNORMAL / BAS / CRITIQUE\nselon quantité vs seuil min"),
    ("vue_commandes_en_attente",  "Toutes les commandes\navec statut PENDING\nprêtes à filtrer"),
    ("vue_delais_fournisseurs",   "Fournisseurs triés par\ndélai de livraison croissant\npour chaque produit"),
]
for i, (name, desc) in enumerate(views):
    x = 0.4 + i * 4.25
    add_rect(s, x, 1.5, 3.9, 0.5, DARK)
    add_text(s, name, x, 1.52, 3.9, 0.46, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, desc, x+0.1, 2.1, 3.7, 1.5, size=13, color=DARK)

code = '''db.create_collection("vue_stock_status", viewOn="products",
  pipeline=[{
    "$project": {
      "_id": 0,
      "nom": "$name",
      "quantite": "$stock.quantity",
      "seuil": "$stock.min_threshold",
      "statut": { "$switch": { "branches": [
          { "case": { "$lt": ["$stock.quantity",
                              "$stock.min_threshold"] },
            "then": "CRITIQUE" },
          { "case": { "$lt": ["$stock.quantity",
             {"$multiply": ["$stock.min_threshold", 2]}] },
            "then": "BAS" }],
        "default": "NORMAL" }}
    }
  }]
)'''
code_block(s, code, 0.4, 3.6, 12.5, 3.6, size=10)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — F1 : État du stock
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "F1 — État du stock d'un produit", "etat_stock(db, product_id)")

add_bullet_box(s, [
    "Requête directe — 0 jointure grâce au stock embarqué dans products",
    "$switch calcule le statut : NORMAL / BAS / CRITIQUE",
    "Seuil BAS = quantité < 2× min_threshold",
], 0.5, 1.5, 6.0, 1.8, size=15, color=DARK)

code = '''pipeline = [
  { "$match": { "product_id": product_id } },
  { "$project": {
      "_id": 0,
      "nom":         "$name",
      "quantite":    "$stock.quantity",
      "seuil_min":   "$stock.min_threshold",
      "statut": { "$switch": { "branches": [
          { "case": { "$lt": ["$stock.quantity",
                              "$stock.min_threshold"] },
            "then": "CRITIQUE — réapprovisionnement urgent" },
          { "case": { "$lt": ["$stock.quantity",
             { "$multiply": ["$stock.min_threshold", 2] }] },
            "then": "BAS — surveiller" }],
        "default": "NORMAL" }}
  }}
]'''
code_block(s, code, 0.4, 3.2, 7.5, 3.9, size=10)

# Résultat exemple
add_text(s, "Résultat exemple", 8.2, 3.2, 4.7, 0.4, size=14, bold=True, color=DARK)
result = '''{
  "nom":      "Ordinateur portable",
  "quantite": 447,
  "seuil_min": 9,
  "statut":   "NORMAL"
}'''
code_block(s, result, 8.2, 3.65, 4.7, 2.5, size=12)

add_rect(s, 8.2, 6.3, 4.7, 0.5, GREEN)
add_text(s, "Démo → /stock sur l'application Flask", 8.2, 6.3, 4.7, 0.5,
         size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — F2 : Commandes bloquées
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "F2 — Commandes bloquées par manque de stock", "commandes_bloquees(db)")

add_bullet_box(s, [
    "$match status=PENDING → $lookup products → $unwind",
    "$expr $gt : quantity_ordered > stock.quantity",
    "Calcule le manque = quantité commandée − stock disponible",
], 0.5, 1.5, 12.0, 1.5, size=15, color=DARK)

code = '''pipeline = [
  { "$match": { "status": "PENDING" } },
  { "$lookup": {
      "from":         "products",
      "localField":   "product_id",
      "foreignField": "product_id",
      "as":           "produit"
  }},
  { "$unwind": "$produit" },
  { "$match": {
      "$expr": { "$gt": ["$quantity_ordered",
                         "$produit.stock.quantity"] }
  }},
  { "$project": {
      "_id": 0,
      "order_id":          1,
      "produit":           "$product_name",
      "quantite_commandee":"$quantity_ordered",
      "stock_disponible":  "$produit.stock.quantity",
      "manque": { "$subtract": ["$quantity_ordered",
                                "$produit.stock.quantity"] }
  }},
  { "$sort": { "date_commande": 1 } }
]'''
code_block(s, code, 0.4, 2.9, 7.5, 4.2, size=10)

result = '''{
  "order_id": 42,
  "produit":  "Clavier mécanique",
  "quantite_commandee": 80,
  "stock_disponible":   23,
  "manque":             57
}'''
add_text(s, "Résultat exemple", 8.2, 2.9, 4.7, 0.4, size=14, bold=True, color=DARK)
code_block(s, result, 8.2, 3.35, 4.7, 2.5, size=12)
add_rect(s, 8.2, 6.3, 4.7, 0.5, GREEN)
add_text(s, "Démo → /commandes sur l'application Flask", 8.2, 6.3, 4.7, 0.5,
         size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — F3 : Délais fournisseurs
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "F3 — Délais de livraison par fournisseur", "delais_fournisseur(db, product_id)")

add_bullet_box(s, [
    "Requête sans jointure — supplier_name et country déjà embarqués",
    "Triée par delivery_days croissant",
    "Affiche aussi le prix unitaire par fournisseur",
], 0.5, 1.5, 12.0, 1.5, size=15, color=DARK)

code = '''pipeline = [
  { "$match": { "product_id": product_id } },
  { "$project": {
      "_id": 0,
      "fournisseur":  "$supplier_name",
      "pays":         "$supplier_country",
      "delai_jours":  "$delivery_days",
      "prix_unitaire":"$unit_price"
  }},
  { "$sort": { "delai_jours": 1 } }
]

# Résultat pour product_id = 1
[
  { "fournisseur": "SpeedLog GmbH", "pays": "Allemagne",
    "delai_jours": 2, "prix_unitaire": 285.00 },
  { "fournisseur": "TechPro Distribution", "pays": "Allemagne",
    "delai_jours": 6, "prix_unitaire": 310.50 },
  { "fournisseur": "GlobalSupply Co", "pays": "Italie",
    "delai_jours": 14, "prix_unitaire": 172.76 }
]'''
code_block(s, code, 0.4, 2.9, 8.5, 4.2, size=10)

add_text(s, "Sans dénormalisation :", 9.2, 2.9, 3.7, 0.4, size=13, bold=True,
         color=RGBColor(0xE7, 0x4C, 0x3C))
add_text(s, "$lookup obligatoire\nsur suppliers\n→ coût supplémentaire", 9.2, 3.35, 3.7, 1.2,
         size=13, color=DARK)
add_text(s, "Avec dénormalisation :", 9.2, 4.7, 3.7, 0.4, size=13, bold=True, color=GREEN)
add_text(s, "supplier_name et country\ndéjà dans le document\n→ 0 jointure", 9.2, 5.15, 3.7, 1.2,
         size=13, color=DARK)
add_rect(s, 9.2, 6.3, 3.7, 0.5, GREEN)
add_text(s, "Démo → /fournisseurs", 9.2, 6.3, 3.7, 0.5, size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — F4 : Date de livraison
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "F4 — Date de livraison d'une commande", "date_livraison(db, order_id)")

add_bullet_box(s, [
    "Si statut = DELIVERED → retourne actual_delivery_date",
    "Sinon → retourne expected_delivery_date + fournisseur le plus rapide",
    "find_one avec sort delivery_days=1 pour trouver le meilleur fournisseur",
], 0.5, 1.5, 12.0, 1.5, size=15, color=DARK)

code = '''order = db.orders.find_one({ "order_id": order_id })

if order["status"] == "DELIVERED":
    return { "livraison_reelle": order["actual_delivery_date"] }

# Commande non livrée : meilleur fournisseur disponible
meilleur = db.supplier_products.find_one(
    { "product_id": order["product_id"] },
    sort=[("delivery_days", 1)]
)
return {
    "statut":              order["status"],
    "livraison_prevue":    order["expected_delivery_date"],
    "fournisseur_rapide":  meilleur["supplier_name"],
    "delai_fournisseur":   meilleur["delivery_days"]
}'''
code_block(s, code, 0.4, 2.9, 7.8, 3.8, size=11)

add_text(s, "Commande livrée", 8.5, 2.9, 4.4, 0.4, size=13, bold=True, color=GREEN)
r1 = '''{ "statut": "DELIVERED",
  "livraison_reelle": "2025-07-11" }'''
code_block(s, r1, 8.5, 3.35, 4.4, 1.0, size=12)

add_text(s, "Commande en attente", 8.5, 4.5, 4.4, 0.4, size=13, bold=True,
         color=RGBColor(0xE7, 0x4C, 0x3C))
r2 = '''{ "statut": "PENDING",
  "livraison_prevue": "2025-09-03",
  "fournisseur_rapide": "SpeedLog GmbH",
  "delai_fournisseur": 2 }'''
code_block(s, r2, 8.5, 4.95, 4.4, 1.5, size=12)
add_rect(s, 8.5, 6.6, 4.4, 0.45, GREEN)
add_text(s, "Démo → /livraison", 8.5, 6.6, 4.4, 0.45, size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — Index & Optimisation
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "Index & Optimisation", "COLLSCAN → IXSCAN")

# Tableau index
h_cols = ["Collection", "Champ(s)", "Type", "Usage"]
rows_idx = [
    ["products",          "product_id",              "Simple",  "F1 — état du stock"],
    ["products",          "stock.quantity",           "Simple",  "Seuil critique"],
    ["products",          "category",                 "Simple",  "Filtrer par catégorie"],
    ["orders",            "status",                   "Simple",  "F2 — PENDING"],
    ["orders",            "(status, product_id)",     "Composé", "F2 optimisé"],
    ["supplier_products", "product_id",               "Simple",  "F3 — délais"],
    ["supplier_products", "(product_id, delivery_days)","Composé","F3 + tri intégré"],
]
col_w = [2.8, 3.0, 1.8, 3.5]
col_x2 = [0.3, 3.15, 6.2, 8.05]
row_y2 = [1.5 + i*0.57 for i in range(len(rows_idx)+1)]

for ci, (hdr, cx, cw) in enumerate(zip(h_cols, col_x2, col_w)):
    add_rect(s, cx, row_y2[0], cw, 0.5, DARK)
    add_text(s, hdr, cx+0.05, row_y2[0]+0.08, cw, 0.35, size=13, bold=True, color=WHITE)

for ri, row in enumerate(rows_idx):
    bg = WHITE if ri % 2 == 0 else RGBColor(0xEB, 0xF5, 0xEB)
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x2, col_w)):
        add_rect(s, cx, row_y2[ri+1], cw, 0.5, bg)
        color = GREEN if ci == 2 and "Composé" in cell else DARK
        add_text(s, cell, cx+0.05, row_y2[ri+1]+0.08, cw, 0.35, size=11, color=color)

add_rect(s, 0.3, 6.45, 5.5, 0.6, RGBColor(0xFF, 0xF3, 0xCD))
add_text(s, "Sans index → COLLSCAN (parcourt tous les documents)", 0.35, 6.5, 5.4, 0.5,
         size=12, color=RGBColor(0x85, 0x62, 0x04))

add_rect(s, 6.2, 6.45, 6.8, 0.6, RGBColor(0xD4, 0xED, 0xDA))
add_text(s, "Avec index → IXSCAN (accès direct via l'index)", 6.25, 6.5, 6.7, 0.5,
         size=12, color=RGBColor(0x15, 0x55, 0x24))


# ═══════════════════════════════════════════════════════════════
# SLIDE 13 — Front Web Flask
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "Front Web Flask", "Architecture & pages fonctionnelles")

# Architecture
add_text(s, "Architecture", 0.4, 1.5, 4.0, 0.4, size=15, bold=True, color=DARK)
add_bullet_box(s, [
    "web/app.py       — routes Flask + PyMongo",
    "web/templates/   — pages HTML (Jinja2)",
    "web/static/      — style CSS",
    "scripts/queries.py — fonctions d'agrégation",
], 0.4, 1.95, 5.5, 2.2, size=13, color=DARK, bullet="  ")

add_text(s, "Lancement", 0.4, 4.2, 4.0, 0.4, size=15, bold=True, color=DARK)
code_block(s, "python web/app.py\n→ http://localhost:5000", 0.4, 4.65, 5.5, 0.9, size=13)

pages = [
    ("/ Dashboard",      "Accueil + stats globales"),
    ("/stock",           "F1 — Recherche par produit"),
    ("/commandes",       "F2 — Commandes bloquées"),
    ("/fournisseurs",    "F3 — Délais fournisseurs"),
    ("/livraison",       "F4 — Suivi d'une commande"),
]
add_text(s, "Pages disponibles", 6.3, 1.5, 6.7, 0.4, size=15, bold=True, color=DARK)
for i, (route, desc) in enumerate(pages):
    y = 1.95 + i * 0.9
    add_rect(s, 6.3, y, 6.7, 0.75, DARK if i == 0 else WHITE)
    col = WHITE if i == 0 else DARK
    add_text(s, route, 6.45, y+0.05, 3.0, 0.35, size=13, bold=True, color=GREEN)
    add_text(s, desc,  9.5,  y+0.05, 3.3, 0.35, size=13, color=col)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14 — Démonstration live
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, DARK)
add_rect(s, 0, 3.3, 13.33, 0.06, GREEN)

add_text(s, "Démonstration", 1, 1.0, 11.33, 1.2,
         size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "live", 1, 2.1, 11.33, 1.0,
         size=48, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

add_text(s, "python web/app.py", 3.5, 4.0, 6.33, 0.7,
         size=24, color=CODE_FG, align=PP_ALIGN.CENTER)
add_text(s, "http://localhost:5000", 3.5, 4.7, 6.33, 0.6,
         size=20, color=GRAY_TEXT, align=PP_ALIGN.CENTER)

add_bullet_box(s, [
    "F1 — État du stock d'un produit",
    "F2 — Commandes bloquées",
    "F3 — Délais fournisseurs",
    "F4 — Date de livraison",
], 4.0, 5.4, 5.33, 2.0, size=16, color=WHITE, bullet="✓  ")


# ═══════════════════════════════════════════════════════════════
# SLIDE 15 — Conclusion
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LIGHT_GRAY)
header(s, "Conclusion", "Bilan du projet")

add_text(s, "Ce qu'on a appris", 0.5, 1.6, 5.5, 0.4, size=16, bold=True, color=DARK)
add_bullet_box(s, [
    "NoSQL vs SQL : quand choisir MongoDB",
    "Dénormalisation : éviter les jointures coûteuses",
    "Pipelines d'agrégation ($match, $lookup, $project, $switch)",
    "Index : COLLSCAN → IXSCAN, mesure avec explain()",
    "Intégration PyMongo + Flask",
], 0.5, 2.05, 5.8, 3.5, size=14, color=DARK)

add_text(s, "Limites", 7.0, 1.6, 5.8, 0.4, size=16, bold=True, color=RGBColor(0xE7, 0x4C, 0x3C))
add_bullet_box(s, [
    "Données mockées (pas de mise à jour temps réel)",
    "Dénormalisation complique les mises à jour",
    "Pas d'authentification sur le front",
], 7.0, 2.05, 5.8, 2.0, size=14, color=DARK)

add_text(s, "Pour aller plus loin", 7.0, 4.1, 5.8, 0.4, size=16, bold=True, color=GREEN)
add_bullet_box(s, [
    "Alertes automatiques si stock critique",
    "Tableau de bord temps réel",
    "API REST complète avec Flask",
], 7.0, 4.5, 5.8, 1.8, size=14, color=DARK)

add_rect(s, 0.5, 6.3, 12.33, 0.7, DARK)
add_text(s, "Merci pour votre attention  •  Kelvya & Maxime  •  Projet C — Suivi des stocks",
         0.5, 6.3, 12.33, 0.7, size=15, color=WHITE, align=PP_ALIGN.CENTER)


# ── Sauvegarde ────────────────────────────────────────────────
output = "presentation_projet_c.pptx"
prs.save(output)
print(f"OK - Presentation generee : {output}")
print(f"   15 slides - ouvre avec PowerPoint ou LibreOffice Impress")
